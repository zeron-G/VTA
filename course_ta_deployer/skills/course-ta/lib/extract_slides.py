#!/usr/bin/env python3
"""
extract_slides.py — PPTX → Compact Markdown index generator for course-ta skill.

Produces lean, summaritive memory files optimized for fast RAG lookup:
- One file per PPTX (not per slide)
- Frontmatter: module, file, slide count, topic list
- Per-slide: title + key bullets only (strips filler, URLs, single words)
- Speaker notes: first 2 sentences only (high-signal context)
- Skips image-only, blank, and boilerplate slides

Usage:
    python3 extract_slides.py <course_dir> <memory_dir>
"""

import os
import sys
import re
import json
import datetime
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
except ImportError:
    print("ERROR: python-pptx not installed. Run: pip3 install python-pptx")
    sys.exit(1)


# ── Text utilities ──────────────────────────────────────────────────────────

BOILERPLATE = re.compile(
    r'^(http[s]?://|www\.|©|copyright|confidential|proprietary'
    r'|all rights reserved|slide \d+|\d+$)',
    re.IGNORECASE
)

def clean_text(t: str) -> str:
    t = t.strip()
    t = re.sub(r'\s+', ' ', t)
    return t

def is_meaningful(line: str) -> bool:
    line = line.strip()
    if len(line) < 4:
        return False
    if BOILERPLATE.search(line):
        return False
    return True

def truncate_notes(notes_text: str, max_sentences: int = 2) -> str:
    """Return first N sentences of speaker notes."""
    if not notes_text:
        return ""
    # Split on sentence boundaries
    sentences = re.split(r'(?<=[.!?])\s+', notes_text.strip())
    chosen = [s.strip() for s in sentences[:max_sentences] if len(s.strip()) > 10]
    return " ".join(chosen)

def slugify(name: str) -> str:
    name = re.sub(r'\.pptx$', '', name, flags=re.IGNORECASE)
    name = name.lower()
    name = re.sub(r'[^a-z0-9]+', '-', name)
    return name.strip('-')


# ── Slide extraction ─────────────────────────────────────────────────────────

def extract_slide(slide, slide_num: int):
    """
    Extract title + key content bullets from one slide.
    Returns None if slide has no useful text (image-only, blank, etc.)
    """
    title = None
    bullets = []
    notes_text = ""

    # Speaker notes
    try:
        if slide.has_notes_slide:
            raw_notes = slide.notes_slide.notes_text_frame.text.strip()
            notes_text = truncate_notes(raw_notes, max_sentences=2)
    except Exception:
        pass

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        try:
            text = shape.text_frame.text.strip()
        except Exception:
            continue
        if not text:
            continue

        # Determine if this is a title placeholder
        is_title = False
        try:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                is_title = True
        except Exception:
            pass

        if is_title:
            title = clean_text(text.splitlines()[0]) if text else None
        else:
            # Collect meaningful lines as bullets
            for line in text.splitlines():
                line = clean_text(line)
                if is_meaningful(line):
                    bullets.append(line)

    if not title and bullets:
        title = bullets.pop(0)

    # Deduplicate bullets while preserving order
    seen = set()
    unique_bullets = []
    for b in bullets:
        key = b.lower()
        if key not in seen:
            seen.add(key)
            unique_bullets.append(b)

    # Skip essentially empty slides
    if not title and not unique_bullets:
        return None

    return {
        "num": slide_num,
        "title": title or f"Slide {slide_num}",
        "bullets": unique_bullets[:12],   # cap at 12 bullets per slide
        "notes": notes_text,
    }


# ── PPTX → Markdown ──────────────────────────────────────────────────────────

def pptx_to_markdown(pptx_path: Path, module_label: str, course_meta: dict) -> str:
    """
    Convert a PPTX to a compact, search-optimised markdown file.

    Structure:
      # <Course> | <Module> | <Lecture title>
      ## Quick Reference  ← topic list for fast triage
      ## Slides           ← per-slide title + bullets + 2-sentence note
    """
    prs = Presentation(str(pptx_path))
    slides_data = []
    for i, slide in enumerate(prs.slides, 1):
        result = extract_slide(slide, i)
        if result:
            slides_data.append(result)

    course_name = course_meta.get("course_name", "Course")
    semester    = course_meta.get("semester", "")
    lecture_id  = pptx_path.stem

    lines = []

    # ── Frontmatter ──
    lines.append(f"# {course_name} | {module_label} | {lecture_id}")
    lines.append(f"> **File:** `{pptx_path.name}`")
    lines.append(f"> **Course:** {course_name} ({semester})")
    lines.append(f"> **Module:** {module_label}")
    lines.append(f"> **Slides:** {len(prs.slides)} total, {len(slides_data)} indexed")
    lines.append(f"> **Indexed:** {datetime.date.today().isoformat()}")
    lines.append("")

    # ── Quick Reference (topic list) ──
    lines.append("## Quick Reference")
    lines.append("")
    lines.append("Key topics covered in this lecture:")
    lines.append("")
    for s in slides_data:
        lines.append(f"- **Slide {s['num']}:** {s['title']}")
    lines.append("")

    # ── Slide Content ──
    lines.append("## Slide Content")
    lines.append("")

    for s in slides_data:
        lines.append(f"### Slide {s['num']}: {s['title']}")
        if s["bullets"]:
            for b in s["bullets"]:
                lines.append(f"- {b}")
        if s["notes"]:
            lines.append(f"> 📝 *{s['notes']}*")
        lines.append("")

    return "\n".join(lines)


# ── Course processing ─────────────────────────────────────────────────────────

def find_course_meta(course_dir: Path) -> dict:
    search = course_dir
    for _ in range(6):
        candidate = search / "course-ta.json"
        if candidate.exists():
            try:
                with open(candidate) as f:
                    return json.load(f)
            except Exception:
                pass
        search = search.parent
    return {}

def process_course(course_dir: Path, memory_dir: Path):
    course_meta = find_course_meta(course_dir)
    course_name = course_dir.name  # the directory name is used as the course slug

    slides_dirs = sorted(course_dir.rglob("slides"))
    if not slides_dirs:
        print(f"  No slides/ directories found under {course_dir}")
        return []

    generated = []
    for slides_dir in slides_dirs:
        module_name = slides_dir.parent.name   # e.g. module1
        module_label = module_name.replace("module", "Module ")  # "Module 1"
        course_meta["module"] = module_label

        pptx_files = sorted(slides_dir.glob("*.pptx"))
        if not pptx_files:
            print(f"  No .pptx files in {slides_dir}")
            continue

        for pptx_path in pptx_files:
            real_path = pptx_path.resolve()
            slug = slugify(pptx_path.name)
            md_name = f"{course_name}__{module_name}__{slug}.md"
            md_path = memory_dir / md_name

            print(f"  [{module_name}] {pptx_path.name} → {md_name}")
            try:
                md_content = pptx_to_markdown(real_path, module_label, course_meta)
                md_path.write_text(md_content, encoding="utf-8")
                generated.append(md_name)
            except Exception as e:
                print(f"    ERROR: {e}")

    print(f"\nGenerated {len(generated)} index files in {memory_dir}")
    return generated


def main():
    if len(sys.argv) < 3:
        print(__doc__)
        sys.exit(1)

    course_dir = Path(sys.argv[1]).resolve()
    memory_dir = Path(sys.argv[2]).resolve()

    if not course_dir.exists():
        print(f"ERROR: course_dir not found: {course_dir}")
        sys.exit(1)
    memory_dir.mkdir(parents=True, exist_ok=True)

    print(f"Course dir:  {course_dir}")
    print(f"Memory dir:  {memory_dir}")
    print("")

    # Auto-detect: if course_dir itself has slides/ subdirs, process it directly;
    # otherwise treat it as a root containing multiple course folders.
    has_slides = any(True for _ in course_dir.rglob("slides"))

    if has_slides:
        process_course(course_dir, memory_dir)
    else:
        for subdir in sorted(d for d in course_dir.iterdir() if d.is_dir()):
            print(f"\n=== Course: {subdir.name} ===")
            process_course(subdir, memory_dir)


if __name__ == "__main__":
    main()
